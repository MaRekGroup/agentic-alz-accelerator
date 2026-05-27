"""
PowerPoint Reporter — generates executive summary slide decks from WARA assessment results.

Produces a branded PowerPoint presentation compatible with Microsoft WARA output:
- Title slide with scope and date
- Assessment overview with methodology and execution summary
- Executive summary with overall score and severity breakdown
- Health dashboard, workload summary, and per-pillar findings
- Impact, remediation, and inventory appendix slides
"""

import logging
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt
except ImportError as exc:
    Presentation = None
    _PPTX_IMPORT_ERROR: ImportError | None = exc
else:
    _PPTX_IMPORT_ERROR = None

from src.tools.discovery import DiscoveryResult
from src.tools.wara_engine import AssessmentResult, Finding

logger = logging.getLogger(__name__)

_PILLAR_NAMES = {
    "security": "Security",
    "reliability": "Reliability",
    "cost_optimization": "Cost Optimization",
    "operational_excellence": "Operational Excellence",
    "performance": "Performance Efficiency",
}

_PILLAR_ORDER = [
    "security",
    "reliability",
    "cost_optimization",
    "operational_excellence",
    "performance",
]

_SEVERITY_ORDER = {"critical": 0, "high": 1, "medium": 2, "low": 3}

_SEVERITY_COLORS = {
    "critical": RGBColor(0xCC, 0x00, 0x00) if RGBColor else None,
    "high": RGBColor(0xE6, 0x73, 0x00) if RGBColor else None,
    "medium": RGBColor(0xCC, 0xA3, 0x00) if RGBColor else None,
    "low": RGBColor(0x33, 0x66, 0xCC) if RGBColor else None,
}

_BRAND_BLUE = RGBColor(0x00, 0x78, 0xD4) if RGBColor else None
_BRAND_DARK = RGBColor(0x24, 0x24, 0x24) if RGBColor else None
_WHITE = RGBColor(0xFF, 0xFF, 0xFF) if RGBColor else None
_LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2) if RGBColor else None
_TEXT_GRAY = RGBColor(0x99, 0x99, 0x99) if RGBColor else None



def _score_color(score: float) -> "RGBColor":
    if score >= 80:
        return RGBColor(0x10, 0x7C, 0x10)
    if score >= 60:
        return RGBColor(0xCA, 0x83, 0x00)
    return RGBColor(0xCC, 0x00, 0x00)


class PptxReporter:
    """Generates PowerPoint executive summary decks for WARA assessments."""

    def __init__(self, output_dir: Path):
        if _PPTX_IMPORT_ERROR is not None:
            raise ImportError("python-pptx is required for PowerPoint report generation") from _PPTX_IMPORT_ERROR
        self.output_dir = Path(output_dir)

    def generate(
        self,
        discovery: DiscoveryResult,
        assessment: AssessmentResult,
        *,
        scope_label: str | None = None,
    ) -> Path:
        """Generate the PowerPoint executive summary deck.

        Args:
            discovery: Discovery inventory for resource context.
            assessment: Assessment results with scores and findings.
            scope_label: Optional label for the title slide.

        Returns:
            Path to the generated .pptx file.
        """
        self.output_dir.mkdir(parents=True, exist_ok=True)
        label = scope_label or discovery.scope

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        self._add_title_slide(prs, label, assessment)
        self._add_assessment_overview(prs, label, assessment)
        self._add_executive_summary(prs, discovery, assessment)
        self._add_health_dashboard(prs, discovery, assessment)
        self._add_workload_summary(prs, discovery)
        for pillar_key in _PILLAR_ORDER:
            self._add_pillar_slide(prs, assessment, pillar_key)
        self._add_impact_slides(prs, assessment.findings, "Critical & High Impact Findings")
        self._add_impact_slides(prs, assessment.findings, "Medium Impact Findings")
        self._add_impact_slides(prs, assessment.findings, "Low Impact Findings")
        self._add_remediation_slide(prs, assessment)
        self._add_inventory_slide(prs, discovery)

        output_path = self.output_dir / "wara-executive-summary.pptx"
        prs.save(str(output_path))
        logger.info("Generated PowerPoint executive summary: %s", output_path)
        return output_path

    def _add_title_slide(self, prs: "Presentation", label: str, assessment: AssessmentResult) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._fill_background(slide, _BRAND_BLUE)

        title_box = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(11), Inches(1.4))
        tf = title_box.text_frame
        tf.word_wrap = True
        title = tf.paragraphs[0]
        title.text = "Well-Architected Reliability Assessment"
        title.font.size = Pt(30)
        title.font.bold = True
        title.font.color.rgb = _WHITE
        title.alignment = PP_ALIGN.LEFT

        subtitle = tf.add_paragraph()
        subtitle.text = "Agentic ALZ Accelerator — Automated Assessment"
        subtitle.font.size = Pt(20)
        subtitle.font.color.rgb = _WHITE
        subtitle.alignment = PP_ALIGN.LEFT
        subtitle.space_before = Pt(8)

        scope_box = slide.shapes.add_textbox(Inches(1), Inches(4.1), Inches(11), Inches(1.3))
        stf = scope_box.text_frame
        stf.word_wrap = True
        scope = stf.paragraphs[0]
        scope.text = f"Scope: {label}"
        scope.font.size = Pt(20)
        scope.font.color.rgb = _WHITE
        scope.alignment = PP_ALIGN.LEFT

        generated = stf.add_paragraph()
        generated.text = f"Generated: {self._format_timestamp(assessment.assessed_at)}"
        generated.font.size = Pt(16)
        generated.font.color.rgb = _WHITE
        generated.alignment = PP_ALIGN.LEFT

    def _add_assessment_overview(self, prs: "Presentation", label: str, assessment: AssessmentResult) -> None:
        slide = self._create_data_slide(prs, "Assessment Overview")
        overview_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.6), Inches(4.8))
        tf = overview_box.text_frame
        tf.word_wrap = True
        bullets = [
            f"Scope: {label}",
            f"Assessment Date: {self._format_timestamp(assessment.assessed_at)}",
            "Methodology: Well-Architected Framework — 5 Pillar Assessment",
            "WAF Pillars: Security, Reliability, Cost Optimization, Operational Excellence, Performance Efficiency",
            f"Checks Executed: {assessment.checks_run}",
            f"Checks Passed: {assessment.checks_passed}",
        ]
        for index, bullet in enumerate(bullets):
            paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = _BRAND_DARK
            paragraph.space_after = Pt(12)
            paragraph.bullet = True

    def _add_executive_summary(
        self,
        prs: "Presentation",
        discovery: DiscoveryResult,
        assessment: AssessmentResult,
    ) -> None:
        slide = self._create_data_slide(prs, "Executive Summary")
        counts = self._severity_counts(assessment.findings)

        score_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4), Inches(2.4))
        tf = score_box.text_frame
        tf.word_wrap = True
        label = tf.paragraphs[0]
        label.text = "Overall Score"
        label.font.size = Pt(16)
        label.font.bold = True
        label.font.color.rgb = _BRAND_DARK

        score = tf.add_paragraph()
        score.text = f"{assessment.overall_score:.0f}/100"
        score.font.size = Pt(44)
        score.font.bold = True
        score.font.color.rgb = _score_color(assessment.overall_score)

        severity = tf.add_paragraph()
        severity.text = (
            f"🔴 {counts['critical']} Critical | 🟠 {counts['high']} High | "
            f"🟡 {counts['medium']} Medium | 🔵 {counts['low']} Low"
        )
        severity.font.size = Pt(12)
        severity.font.color.rgb = _BRAND_DARK
        severity.space_before = Pt(4)

        stats_box = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(3.8), Inches(1.8))
        stf = stats_box.text_frame
        stf.word_wrap = True
        for index, stat in enumerate(self._summary_stats(discovery, assessment)):
            paragraph = stf.paragraphs[0] if index == 0 else stf.add_paragraph()
            paragraph.text = stat
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = _BRAND_DARK
            paragraph.space_after = Pt(4)

        self._add_pillar_score_table(slide, assessment, left=Inches(5.1), top=Inches(1.5))

    def _add_health_dashboard(
        self,
        prs: "Presentation",
        discovery: DiscoveryResult,
        assessment: AssessmentResult,
    ) -> None:
        slide = self._create_data_slide(prs, "Health & Risk Dashboard")
        counts = self._severity_counts(assessment.findings)
        metrics = [
            ("Total Findings", str(len(assessment.findings)), _BRAND_DARK, 30),
            ("Total Resources", str(self._resource_total(discovery)), _BRAND_DARK, 30),
            ("Critical", str(counts["critical"]), _SEVERITY_COLORS["critical"], 24),
            ("High", str(counts["high"]), _SEVERITY_COLORS["high"], 24),
            ("Medium", str(counts["medium"]), _SEVERITY_COLORS["medium"], 24),
            ("Low", str(counts["low"]), _SEVERITY_COLORS["low"], 24),
        ]

        for index, (label, value, color, size) in enumerate(metrics):
            top = Inches(1.4 + index * 0.8)
            self._add_metric_box(slide, label, value, Inches(0.8), top, Inches(3.4), color, size)

        impacted = self._top_impacted_services(assessment.findings)
        table_shape = slide.shapes.add_table(7, 3, Inches(4.6), Inches(1.5), Inches(7.8), Inches(3.6))
        table = table_shape.table
        headers = ["Resource Type", "Finding Count", "Top Severity"]
        for column, header in enumerate(headers):
            cell = table.cell(0, column)
            cell.text = header
            self._style_header_cell(cell)

        for row in range(1, 7):
            if row <= len(impacted):
                resource_type, count, severity = impacted[row - 1]
                values = [resource_type, str(count), severity.capitalize()]
            else:
                values = ["—", "0", "—"]
            for column, value in enumerate(values):
                cell = table.cell(row, column)
                cell.text = value
                self._style_body_cell(cell, 11)
            severity_key = values[2].lower()
            if severity_key in _SEVERITY_COLORS:
                self._style_severity_cell(table.cell(row, 2), severity_key)

        table.columns[0].width = Inches(4.8)
        table.columns[1].width = Inches(1.4)
        table.columns[2].width = Inches(1.6)

        label_box = slide.shapes.add_textbox(Inches(4.6), Inches(1.1), Inches(4), Inches(0.3))
        label = label_box.text_frame.paragraphs[0]
        label.text = "Top Impacted Services"
        label.font.size = Pt(16)
        label.font.bold = True
        label.font.color.rgb = _BRAND_DARK

    def _add_workload_summary(self, prs: "Presentation", discovery: DiscoveryResult) -> None:
        slide = self._create_data_slide(prs, "Resource Inventory Summary")
        resource_rows = self._resource_rows(discovery)[:15]
        if resource_rows:
            self._add_resource_table(slide, resource_rows, Inches(0.8), Inches(1.5), Inches(8.5))
        else:
            self._add_message_box(slide, "No resources discovered.", Inches(1), Inches(2), Inches(10))

        total_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(4), Inches(0.4))
        paragraph = total_box.text_frame.paragraphs[0]
        paragraph.text = f"Total resources: {self._resource_total(discovery)}"
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = _BRAND_DARK

    def _add_pillar_slide(self, prs: "Presentation", assessment: AssessmentResult, pillar_key: str) -> None:
        slide = self._create_data_slide(prs, f"Pillar: {_PILLAR_NAMES.get(pillar_key, pillar_key)}")
        pillar_score = assessment.pillar_scores.get(pillar_key)
        score = pillar_score.score if pillar_score else 100.0

        score_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(3.4), Inches(1.2))
        tf = score_box.text_frame
        tf.word_wrap = True
        score_paragraph = tf.paragraphs[0]
        score_paragraph.text = f"{score:.0f}/100"
        score_paragraph.font.size = Pt(44)
        score_paragraph.font.bold = True
        score_paragraph.font.color.rgb = _score_color(score)

        counts_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(3.4), Inches(2.1))
        ctf = counts_box.text_frame
        ctf.word_wrap = True
        counts = self._pillar_counts(pillar_score)
        for index, count in enumerate(counts):
            paragraph = ctf.paragraphs[0] if index == 0 else ctf.add_paragraph()
            paragraph.text = count
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = _BRAND_DARK
            paragraph.space_after = Pt(3)

        pillar_findings = [finding for finding in assessment.findings if finding.pillar == pillar_key]
        top_findings = sorted(pillar_findings, key=lambda finding: _SEVERITY_ORDER.get(finding.severity, 99))[:5]
        if top_findings:
            self._add_findings_table(slide, top_findings, left=Inches(4.6), top=Inches(1.5))
        else:
            self._add_message_box(slide, "No findings for this pillar.", Inches(4.9), Inches(2.2), Inches(6.5))

    def _add_impact_slides(
        self,
        prs: "Presentation",
        findings: list[Finding],
        title: str,
        max_rows: int = 10,
    ) -> int:
        severities = self._impact_severities(title)
        filtered = [finding for finding in findings if finding.severity in severities]
        chunks = [filtered[index : index + max_rows] for index in range(0, len(filtered), max_rows)] or [[]]

        for index, chunk in enumerate(chunks, start=1):
            slide_title = title if len(chunks) == 1 else f"{title} ({index}/{len(chunks)})"
            slide = self._create_data_slide(prs, slide_title)
            if not chunk:
                self._add_message_box(slide, "No findings.", Inches(1), Inches(2), Inches(10))
                continue

            table_shape = slide.shapes.add_table(len(chunk) + 1, 5, Inches(0.4), Inches(1.4), Inches(12.4), Inches(5.3))
            table = table_shape.table
            headers = ["#", "Rule ID", "Severity", "Title", "Recommendation"]
            for column, header in enumerate(headers):
                cell = table.cell(0, column)
                cell.text = header
                self._style_header_cell(cell)

            for row, finding in enumerate(chunk, start=1):
                values = [
                    str(row + ((index - 1) * max_rows)),
                    finding.rule_id,
                    finding.severity.capitalize(),
                    finding.title,
                    finding.recommendation,
                ]
                for column, value in enumerate(values):
                    cell = table.cell(row, column)
                    cell.text = value
                    self._style_body_cell(cell, 10)
                self._style_severity_cell(table.cell(row, 2), finding.severity)

            table.columns[0].width = Inches(0.6)
            table.columns[1].width = Inches(1.3)
            table.columns[2].width = Inches(1.2)
            table.columns[3].width = Inches(3.3)
            table.columns[4].width = Inches(6)
        return len(chunks)

    def _add_remediation_slide(self, prs: "Presentation", assessment: AssessmentResult) -> None:
        slide = self._create_data_slide(prs, "Remediation Roadmap")
        top_findings = sorted(
            assessment.findings,
            key=lambda finding: (_SEVERITY_ORDER.get(finding.severity, 99), finding.rule_id),
        )[:10]
        if not top_findings:
            self._add_message_box(slide, "No findings — all checks passed.", Inches(1), Inches(2), Inches(10))
            return

        table_shape = slide.shapes.add_table(len(top_findings) + 1, 5, Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        table = table_shape.table
        headers = ["Priority", "Rule ID", "Severity", "Title", "Recommendation"]
        for column, header in enumerate(headers):
            cell = table.cell(0, column)
            cell.text = header
            self._style_header_cell(cell)

        for row, finding in enumerate(top_findings, start=1):
            values = [str(row), finding.rule_id, finding.severity.capitalize(), finding.title, finding.recommendation]
            for column, value in enumerate(values):
                cell = table.cell(row, column)
                cell.text = value
                self._style_body_cell(cell, 10)
            self._style_severity_cell(table.cell(row, 2), finding.severity)

        table.columns[0].width = Inches(0.8)
        table.columns[1].width = Inches(1.2)
        table.columns[2].width = Inches(1.2)
        table.columns[3].width = Inches(4)
        table.columns[4].width = Inches(4.8)

    def _add_inventory_slide(self, prs: "Presentation", discovery: DiscoveryResult) -> None:
        slide = self._create_data_slide(prs, "Appendix: Full Resource Inventory")
        resource_rows = self._resource_rows(discovery)[:15]
        if resource_rows:
            self._add_resource_table(slide, resource_rows, Inches(1), Inches(1.5), Inches(8))
        else:
            self._add_message_box(slide, "No resources discovered.", Inches(1), Inches(2), Inches(10))

    def _create_data_slide(self, prs: "Presentation", title: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._fill_background(slide, _WHITE)
        self._add_slide_title(slide, title)
        self._add_footer(slide)
        return slide

    def _add_pillar_score_table(self, slide, assessment: AssessmentResult, left: "Emu", top: "Emu") -> None:
        table_shape = slide.shapes.add_table(len(_PILLAR_ORDER) + 1, 4, left, top, Inches(7.5), Inches(2.4))
        table = table_shape.table
        headers = ["Pillar", "Score", "Findings", "Critical"]
        for column, header in enumerate(headers):
            cell = table.cell(0, column)
            cell.text = header
            self._style_header_cell(cell)

        for row, pillar_key in enumerate(_PILLAR_ORDER, start=1):
            pillar_score = assessment.pillar_scores.get(pillar_key)
            score = pillar_score.score if pillar_score else 100.0
            findings = pillar_score.findings_count if pillar_score else 0
            critical = pillar_score.critical if pillar_score else 0
            values = [_PILLAR_NAMES[pillar_key], f"{score:.0f}", str(findings), str(critical)]
            for column, value in enumerate(values):
                cell = table.cell(row, column)
                cell.text = value
                self._style_body_cell(cell, 12)
            score_cell = table.cell(row, 1)
            score_cell.fill.solid()
            score_cell.fill.fore_color.rgb = _score_color(score)
            for paragraph in score_cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.color.rgb = _WHITE

        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(1.5)
        table.columns[2].width = Inches(1.5)
        table.columns[3].width = Inches(1.5)

    def _add_findings_table(self, slide, findings: list[Finding], left: "Emu", top: "Emu") -> None:
        table_shape = slide.shapes.add_table(len(findings) + 1, 3, left, top, Inches(7.6), Inches(2.6))
        table = table_shape.table
        headers = ["Severity", "Rule ID", "Title"]
        for column, header in enumerate(headers):
            cell = table.cell(0, column)
            cell.text = header
            self._style_header_cell(cell)

        for row, finding in enumerate(findings, start=1):
            values = [finding.severity.capitalize(), finding.rule_id, finding.title]
            for column, value in enumerate(values):
                cell = table.cell(row, column)
                cell.text = value
                self._style_body_cell(cell, 11)
            self._style_severity_cell(table.cell(row, 0), finding.severity)

        table.columns[0].width = Inches(1.2)
        table.columns[1].width = Inches(1.5)
        table.columns[2].width = Inches(4.9)

    def _add_resource_table(
        self,
        slide,
        resource_rows: list[tuple[str, int]],
        left: "Emu",
        top: "Emu",
        width: "Emu",
    ) -> None:
        table_height = Inches(0.3 * (len(resource_rows) + 1))
        table_shape = slide.shapes.add_table(len(resource_rows) + 1, 2, left, top, width, table_height)
        table = table_shape.table
        for column, header in enumerate(["Resource Type", "Count"]):
            cell = table.cell(0, column)
            cell.text = header
            self._style_header_cell(cell)

        for row, (resource_type, count) in enumerate(resource_rows, start=1):
            for column, value in enumerate([resource_type, str(count)]):
                cell = table.cell(row, column)
                cell.text = value
                self._style_body_cell(cell, 11)

        table.columns[0].width = Inches(6)
        table.columns[1].width = Inches(2)

    def _add_metric_box(
        self,
        slide,
        label: str,
        value: str,
        left: "Emu",
        top: "Emu",
        width: "Emu",
        color: "RGBColor",
        value_size: int,
    ) -> None:
        box = slide.shapes.add_textbox(left, top, width, Inches(0.7))
        tf = box.text_frame
        tf.word_wrap = True
        title = tf.paragraphs[0]
        title.text = label
        title.font.size = Pt(12)
        title.font.color.rgb = _BRAND_DARK

        metric = tf.add_paragraph()
        metric.text = value
        metric.font.size = Pt(value_size)
        metric.font.bold = True
        metric.font.color.rgb = color
        metric.space_before = Pt(2)

    def _add_message_box(self, slide, text: str, left: "Emu", top: "Emu", width: "Emu") -> None:
        message_box = slide.shapes.add_textbox(left, top, width, Inches(1))
        paragraph = message_box.text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = _BRAND_DARK

    def _add_slide_title(self, slide, text: str) -> None:
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        paragraph = title_box.text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True
        paragraph.font.color.rgb = _BRAND_DARK
        paragraph.alignment = PP_ALIGN.LEFT

    def _add_footer(self, slide) -> None:
        footer_box = slide.shapes.add_textbox(Inches(9.5), Inches(7.0), Inches(3.2), Inches(0.25))
        paragraph = footer_box.text_frame.paragraphs[0]
        paragraph.text = "AUTOMATICALLY GENERATED — Please Review"
        paragraph.font.name = "Segoe UI"
        paragraph.font.size = Pt(8)
        paragraph.font.color.rgb = _TEXT_GRAY
        paragraph.alignment = PP_ALIGN.RIGHT

    @staticmethod
    def _fill_background(slide, color: "RGBColor") -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    @staticmethod
    def _style_header_cell(cell) -> None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = _BRAND_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = _WHITE
            paragraph.alignment = PP_ALIGN.LEFT

    @staticmethod
    def _style_body_cell(cell, font_size: int) -> None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = _WHITE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = _BRAND_DARK
            paragraph.alignment = PP_ALIGN.LEFT

    @staticmethod
    def _style_severity_cell(cell, severity: str) -> None:
        color = _SEVERITY_COLORS.get(severity)
        if color is None:
            return
        cell.fill.solid()
        cell.fill.fore_color.rgb = color
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = _WHITE
            paragraph.font.bold = True

    @staticmethod
    def _summary_stats(discovery: DiscoveryResult, assessment: AssessmentResult) -> list[str]:
        return [
            f"Resources: {PptxReporter._resource_total(discovery)}",
            f"Checks Run: {assessment.checks_run}",
            f"Checks Passed: {assessment.checks_passed}",
            f"Findings: {len(assessment.findings)}",
        ]

    @staticmethod
    def _severity_counts(findings: list[Finding]) -> Counter[str]:
        counts: Counter[str] = Counter(dict.fromkeys(_SEVERITY_ORDER, 0))
        counts.update(finding.severity for finding in findings)
        return counts

    @staticmethod
    def _pillar_counts(pillar_score) -> list[str]:
        if not pillar_score or pillar_score.findings_count == 0:
            return ["No findings"]
        return [
            f"Critical: {pillar_score.critical}",
            f"High: {pillar_score.high}",
            f"Medium: {pillar_score.medium}",
            f"Low: {pillar_score.low}",
            f"Total: {pillar_score.findings_count}",
        ]

    @staticmethod
    def _impact_severities(title: str) -> set[str]:
        if title == "Critical & High Impact Findings":
            return {"critical", "high"}
        if title == "Medium Impact Findings":
            return {"medium"}
        return {"low"}

    @staticmethod
    def _finding_groups(finding: Finding) -> list[str]:
        groups: dict[str, str] = {}
        for evidence in finding.evidence:
            if not isinstance(evidence, dict):
                continue
            resource_type = evidence.get("resource_type") or evidence.get("resourceType") or evidence.get("type")
            if isinstance(resource_type, str) and resource_type:
                groups.setdefault(resource_type.lower(), resource_type)
        if groups:
            return list(groups.values())
        pillar_name = _PILLAR_NAMES.get(finding.pillar, finding.pillar.replace("_", " ").title())
        return [pillar_name]

    @staticmethod
    def _top_impacted_services(findings: list[Finding]) -> list[tuple[str, int, str]]:
        counts: Counter[str] = Counter()
        top_severity: dict[str, str] = {}
        display_names: dict[str, str] = {}
        for finding in findings:
            for group in PptxReporter._finding_groups(finding):
                key = group.lower()
                counts[key] += 1
                display_names.setdefault(key, group)
                current = top_severity.get(key)
                if current is None or _SEVERITY_ORDER.get(finding.severity, 99) < _SEVERITY_ORDER.get(current, 99):
                    top_severity[key] = finding.severity
        rows = [
            (display_names[key], count, top_severity.get(key, "low"))
            for key, count in counts.items()
        ]
        rows.sort(key=lambda row: (-row[1], _SEVERITY_ORDER.get(row[2], 99), row[0]))
        return rows[:6]

    @staticmethod
    def _format_timestamp(raw_value: str) -> str:
        if not raw_value:
            return datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
        try:
            dt = datetime.fromisoformat(raw_value.replace("Z", "+00:00"))
        except ValueError:
            return raw_value
        return dt.astimezone(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")

    @staticmethod
    def _resource_rows(discovery: DiscoveryResult) -> list[tuple[str, int]]:
        by_type = discovery.resources.get("by_type", {}) if isinstance(discovery.resources, dict) else {}
        if isinstance(by_type, dict):
            rows = [(resource_type, int(count)) for resource_type, count in by_type.items()]
        else:
            rows = [
                (item.get("type", "unknown"), int(item.get("resource_count", item.get("count", 0))))
                for item in by_type
                if isinstance(item, dict)
            ]
        return sorted(rows, key=lambda row: row[1], reverse=True)

    @staticmethod
    def _resource_total(discovery: DiscoveryResult) -> int:
        if not isinstance(discovery.resources, dict):
            return 0
        total = discovery.resources.get("total_count", discovery.resources.get("total"))
        if total is not None:
            return int(total)
        return sum(count for _, count in PptxReporter._resource_rows(discovery))
