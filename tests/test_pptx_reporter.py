"""Tests for the PowerPoint reporter."""

from copy import deepcopy

import pytest
from pptx import Presentation

from src.tools.discovery import DiscoveryResult, DiscoveryScope
from src.tools.pptx_reporter import PptxReporter
from src.tools.wara_engine import AssessmentResult, Finding, PillarScore


@pytest.fixture
def sample_discovery():
    return DiscoveryResult(
        scope="test-sub",
        scope_type=DiscoveryScope.SUBSCRIPTION,
        resources={
            "by_type": [
                {"type": "microsoft.compute/virtualmachines", "count": 5},
                {"type": "microsoft.storage/storageaccounts", "count": 2},
            ],
            "total_count": 7,
        },
    )


@pytest.fixture
def sample_assessment():
    findings = [
        Finding(
            rule_id="SEC-001",
            title="TLS check",
            pillar="security",
            caf_area="security",
            alz_area="security",
            severity="critical",
            confidence="high",
            recommendation="Fix TLS",
            evidence=[{"resourceType": "microsoft.compute/virtualmachines"}],
            remediation_steps=["Step 1"],
        ),
        Finding(
            rule_id="REL-001",
            title="Backup check",
            pillar="reliability",
            caf_area="management",
            alz_area="management",
            severity="high",
            confidence="high",
            recommendation="Enable backup",
            evidence=[{"type": "microsoft.storage/storageaccounts"}],
            remediation_steps=["Enable backup"],
        ),
    ]
    result = AssessmentResult(scope="test-sub")
    result.findings = findings
    result.checks_run = 10
    result.checks_passed = 8
    result.overall_score = 85.0
    result.pillar_scores = {
        "security": PillarScore(pillar="security", score=80.0, findings_count=1, critical=1),
        "reliability": PillarScore(pillar="reliability", score=90.0, findings_count=1, high=1),
        "cost_optimization": PillarScore(pillar="cost_optimization", score=100.0),
        "operational_excellence": PillarScore(pillar="operational_excellence", score=100.0),
        "performance": PillarScore(pillar="performance", score=100.0),
    }
    return result


@pytest.fixture
def reporter(tmp_path):
    return PptxReporter(output_dir=tmp_path)



def _load_presentation(reporter, sample_discovery, sample_assessment, scope_label="test-sub"):
    output_path = reporter.generate(sample_discovery, sample_assessment, scope_label=scope_label)
    return Presentation(str(output_path))



def _slide_text(slide) -> str:
    return "\n".join(shape.text_frame.text for shape in slide.shapes if shape.has_text_frame)



def _find_slide_index(prs: Presentation, title: str) -> int:
    for index, slide in enumerate(prs.slides):
        if title in _slide_text(slide):
            return index
    raise AssertionError(f"Slide titled '{title}' not found")



def test_generate_creates_pptx_file(reporter, sample_discovery, sample_assessment):
    output_path = reporter.generate(sample_discovery, sample_assessment, scope_label="test-sub")

    assert output_path.exists()
    assert output_path.suffix == ".pptx"
    assert output_path.name == "wara-executive-summary.pptx"



def test_slide_count(reporter, sample_discovery, sample_assessment):
    prs = _load_presentation(reporter, sample_discovery, sample_assessment)
    assert len(prs.slides) == 15



def test_title_slide_content(reporter, sample_discovery, sample_assessment):
    prs = _load_presentation(reporter, sample_discovery, sample_assessment)
    combined = _slide_text(prs.slides[0])
    assert "Well-Architected Reliability Assessment" in combined
    assert "Agentic ALZ Accelerator — Automated Assessment" in combined
    assert "test-sub" in combined



def test_overview_slide_content(reporter, sample_discovery, sample_assessment):
    prs = _load_presentation(reporter, sample_discovery, sample_assessment)
    combined = _slide_text(prs.slides[1])
    assert "Assessment Overview" in combined
    assert "Methodology: Well-Architected Framework — 5 Pillar Assessment" in combined
    assert "Checks Executed: 10" in combined



def test_executive_summary_shows_score(reporter, sample_discovery, sample_assessment):
    prs = _load_presentation(reporter, sample_discovery, sample_assessment)
    combined = _slide_text(prs.slides[2])
    assert "85/100" in combined
    assert "1 Critical" in combined
    assert "1 High" in combined



def test_health_dashboard_slide(reporter, sample_discovery, sample_assessment):
    prs = _load_presentation(reporter, sample_discovery, sample_assessment)
    combined = _slide_text(prs.slides[3])
    assert "Health & Risk Dashboard" in combined
    assert "Total Findings" in combined
    assert "Critical" in combined
    assert any(shape.has_table for shape in prs.slides[3].shapes)



def test_workload_summary_slide(reporter, sample_discovery, sample_assessment):
    prs = _load_presentation(reporter, sample_discovery, sample_assessment)
    combined = _slide_text(prs.slides[4])
    assert "Resource Inventory Summary" in combined
    assert "Total resources: 7" in combined
    assert any(shape.has_table for shape in prs.slides[4].shapes)



def test_remediation_slide_has_findings(reporter, sample_discovery, sample_assessment):
    prs = _load_presentation(reporter, sample_discovery, sample_assessment)
    remediation_slide = prs.slides[13]
    assert "Remediation Roadmap" in _slide_text(remediation_slide)
    assert any(shape.has_table for shape in remediation_slide.shapes)



def test_inventory_slide_has_resources(reporter, sample_discovery, sample_assessment):
    prs = _load_presentation(reporter, sample_discovery, sample_assessment)
    inventory_slide = prs.slides[14]
    assert "Appendix: Full Resource Inventory" in _slide_text(inventory_slide)
    assert any(shape.has_table for shape in inventory_slide.shapes)



def test_impact_slides_auto_duplicate(reporter, sample_discovery, sample_assessment):
    assessment = deepcopy(sample_assessment)
    assessment.findings = [
        Finding(
            rule_id=f"REL-{index:03d}",
            title=f"Impact finding {index}",
            pillar="reliability",
            caf_area="management",
            alz_area="management",
            severity="high",
            confidence="high",
            recommendation="Review impact",
            evidence=[{"type": "microsoft.compute/virtualmachines"}],
        )
        for index in range(1, 16)
    ]
    assessment.pillar_scores["reliability"] = PillarScore(pillar="reliability", score=0.0, findings_count=15, high=15)

    output_path = reporter.generate(sample_discovery, assessment, scope_label="test-sub")
    prs = Presentation(str(output_path))

    impact_titles = [
        _slide_text(slide)
        for slide in prs.slides
        if "Critical & High Impact Findings" in _slide_text(slide)
    ]
    assert len(impact_titles) == 2
    assert len(prs.slides) == 16



def test_auto_generated_footer(reporter, sample_discovery, sample_assessment):
    prs = _load_presentation(reporter, sample_discovery, sample_assessment)
    slides = list(prs.slides)
    assert "AUTOMATICALLY GENERATED — Please Review" not in _slide_text(slides[0])
    for slide in slides[1:]:
        assert "AUTOMATICALLY GENERATED — Please Review" in _slide_text(slide)



def test_empty_assessment(reporter, sample_discovery):
    empty_assessment = AssessmentResult(scope="test-sub")
    empty_assessment.pillar_scores = {
        "security": PillarScore(pillar="security"),
        "reliability": PillarScore(pillar="reliability"),
        "cost_optimization": PillarScore(pillar="cost_optimization"),
        "operational_excellence": PillarScore(pillar="operational_excellence"),
        "performance": PillarScore(pillar="performance"),
    }
    output_path = reporter.generate(sample_discovery, empty_assessment, scope_label="empty")
    prs = Presentation(str(output_path))
    assert len(prs.slides) == 15
    assert "No findings." in _slide_text(prs.slides[10])
    assert "No findings." in _slide_text(prs.slides[11])
    assert "No findings." in _slide_text(prs.slides[12])



def test_generate_all_includes_pptx(tmp_path, sample_discovery, sample_assessment):
    from src.tools.report_generator import ReportGenerator

    gen = ReportGenerator(output_dir=tmp_path)
    outputs = gen.generate_all(sample_discovery, sample_assessment, scope_label="test")
    assert "pptx_executive_summary" in outputs
    assert outputs["pptx_executive_summary"].suffix == ".pptx"
