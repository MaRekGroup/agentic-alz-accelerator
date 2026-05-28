"""Shared utilities for dark-theme Microsoft-branded PowerPoint generation."""

from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu

# Theme colors - dark
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_NAVY = RGBColor(0x00, 0x1B, 0x44)
DARK_BLUE = RGBColor(0x00, 0x20, 0x50)
CARD_BG = RGBColor(0x00, 0x2B, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEXT = RGBColor(0xD9, 0xE8, 0xF5)
SUBTLE_TEXT = RGBColor(0xA0, 0xC4, 0xE8)
SUCCESS_GREEN = RGBColor(0x10, 0x7C, 0x10)
WARNING_AMBER = RGBColor(0xFF, 0xB9, 0x00)
FONT_NAME = "Segoe UI"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

_THEME = 'dark'

def set_theme(theme: str):
    global _THEME
    _THEME = theme

def create_presentation():
    prs = Presentation()
    prs.slide_width = int(SLIDE_WIDTH)
    prs.slide_height = int(SLIDE_HEIGHT)
    return prs

def save_presentation(prs, path):
    Path(path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    print(f"Saved: {path} ({Path(path).stat().st_size} bytes)")

def add_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(layout)

def add_background(slide, color=None):
    if color is None:
        color = DARK_NAVY
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, int(SLIDE_HEIGHT - Inches(0.16)), int(SLIDE_WIDTH), int(Inches(0.16))
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

def style_run(run, font_size, color=None, bold=False):
    if color is None:
        color = WHITE
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color

def add_title(slide, title, font_size=34):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.7))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    style_run(r, font_size, bold=True)
    p.alignment = PP_ALIGN.LEFT

def add_subtitle(slide, text, top=1.1):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    style_run(r, 18, color=SUBTLE_TEXT)

def add_notes(slide, notes):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = notes
    style_run(r, 12, color=RGBColor(0x00, 0x00, 0x00))

def _enable_bullet(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for bu_none in pPr.findall("a:buNone", nsmap):
        pPr.remove(bu_none)
    bu_char = etree.SubElement(pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar")
    bu_char.set("char", "\u2022")
    bu_sz = etree.SubElement(pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buSzPct")
    bu_sz.set("val", "100000")

def add_bullet_list(slide, items, left=0.8, top=1.7, width=11.5, height=5.2, font_size=20, color=None):
    if color is None:
        color = LIGHT_TEXT
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(8)
        p.space_after = Pt(4)
        _enable_bullet(p)
        r = p.add_run()
        r.text = item
        style_run(r, font_size, color=color)

def add_two_column_bullets(slide, left_items, right_items, top=1.7, font_size=18):
    add_bullet_list(slide, left_items, left=0.6, top=top, width=5.8, font_size=font_size)
    add_bullet_list(slide, right_items, left=6.8, top=top, width=5.8, font_size=font_size)

def add_table(slide, headers, rows, left=0.6, top=1.8, width=12.1, row_height=0.45):
    col_count = len(headers)
    row_count = len(rows) + 1
    height = row_height * row_count
    tbl_shape = slide.shapes.add_table(
        row_count, col_count, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    table = tbl_shape.table
    col_width = width / col_count
    for i in range(col_count):
        table.columns[i].width = Inches(col_width)
    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.name = FONT_NAME
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE if ri % 2 == 0 else CARD_BG

def add_card(slide, text, left, top, width=3.8, height=1.2, bg=None):
    if bg is None:
        bg = CARD_BG
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    style_run(r, 14, color=WHITE)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
