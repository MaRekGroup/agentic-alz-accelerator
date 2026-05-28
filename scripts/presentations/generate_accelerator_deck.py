"""Generate a Microsoft-branded PowerPoint deck for the Agentic Azure Enterprise Landing Zone Accelerator."""

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_NAVY = RGBColor(0x00, 0x1B, 0x44)
DARK_BLUE = RGBColor(0x00, 0x20, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEXT = RGBColor(0xD9, 0xE8, 0xF5)
FONT_NAME = "Segoe UI"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
OUTPUT_PATH = Path(__file__).resolve().parent.parent.parent / "docs" / "Agentic_Azure_EALZ_Accelerator_Deck.pptx"


SLIDE_7_ROWS = [
    ["🧠 Conductor", "Master orchestration", "—"],
    ["🔍 Assessor", "Brownfield WAF assessment", "0"],
    ["📜 Scribe", "Requirements capture (8 CAF areas)", "1"],
    ["🏛️ Oracle", "WAF assessment + cost estimation", "2"],
    ["🛡️ Warden", "Policy discovery + security baseline", "3.5"],
    ["📐 Strategist", "AVM module selection + planning", "4"],
    ["⚒️ Forge", "Bicep/Terraform generation (AVM-first)", "5"],
    ["🚀 Envoy", "Deploy with what-if preview", "6"],
    ["🔭 Sentinel", "Continuous compliance monitoring", "8"],
    ["🔧 Mender", "Auto-remediation with rollback", "9"],
    ["⚔️ Challenger", "Adversarial review at gates", "Gates"],
]

SLIDE_8_ROWS = [
    ["Billing & Tenant", "Management group architecture"],
    ["Identity & Access", "Entra ID, RBAC, Managed Identity"],
    ["Resource Organization", "Governance policies, naming"],
    ["Network Topology", "Hub-spoke, Virtual WAN, Private Link"],
    ["Security", "Defender, Key Vault, TLS 1.2"],
    ["Management", "Monitor, Log Analytics, Automation"],
    ["Governance", "Policy, compliance, cost controls"],
    ["Platform Automation", "GitHub Actions CI/CD with OIDC"],
]

SLIDE_11_ROWS = [
    ["Azure Pricing", "Real-time cost estimation", "18 tools"],
    ["Azure Platform", "Live Azure operations", "27 tools"],
    ["Draw.io", "Diagram generation", "Azure icons"],
]

SLIDE_PROFILES_ROWS = [
    ["platform-management", "Logging, monitoring, automation", "1"],
    ["platform-connectivity", "Hub networking, Firewall, DNS", "2"],
    ["platform-identity", "AD DS, PIM, RBAC", "3"],
    ["platform-security", "Sentinel, Defender, SOAR", "4"],
    ["online", "Internet-facing apps, WAF", "App"],
    ["corp", "Internal, private connectivity", "App"],
    ["sandbox", "Dev/test, relaxed policies", "App"],
    ["sap", "SAP workloads, ANF, proximity", "App"],
]

SLIDE_PIPELINE_ROWS = [
    ["Resolve", "Validate Azure access via OIDC"],
    ["Validate", "Security baseline + cost governance + lint"],
    ["Plan", "Bicep what-if / Terraform plan"],
    ["Deploy", "Apply with timestamped deployment name"],
    ["Verify", "Compliance scan + TDD generation"],
]

SLIDE_TDD_ROWS = [
    ["Cover Page", "Deployment metadata, subscription, profile"],
    ["Architecture Diagram", "SVG with Azure icons → PNG"],
    ["Resource Inventory", "Full listing from Resource Graph"],
    ["Network Topology", "VNets, subnets, peering, DNS"],
    ["Security Posture", "6 baseline rules, Defender plans, Policy"],
    ["Compliance Status", "Live scan results post-deploy"],
    ["Cost Governance", "Budgets, alerts, tag compliance"],
    ["Operational Model", "Monitoring, DR, change management"],
]


def style_run(run, font_size: int, color: RGBColor = WHITE, bold: bool = False) -> None:
    """Apply consistent font styling to a run."""
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_background(slide, color: RGBColor = DARK_NAVY) -> None:
    """Apply a branded background and accent bar."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, int(SLIDE_HEIGHT - Inches(0.16)), int(SLIDE_WIDTH), int(Inches(0.16))
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()


def add_title(slide, title: str) -> None:
    """Add a slide title."""
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.7))
    text_frame = box.text_frame
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    style_run(run, 34, bold=True)
    paragraph.alignment = PP_ALIGN.LEFT


def add_notes(slide, notes: str) -> None:
    """Populate the speaker notes for a slide."""
    text_frame = slide.notes_slide.notes_text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = notes
    style_run(run, 12, color=RGBColor(0x00, 0x00, 0x00))


def _enable_bullet(paragraph) -> None:
    """Enable a filled-circle bullet on a paragraph via Open XML."""
    pPr = paragraph._p.get_or_add_pPr()
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    # Remove any existing buNone (no-bullet marker)
    for bu_none in pPr.findall("a:buNone", nsmap):
        pPr.remove(bu_none)
    # Add bullet character
    bu_char = etree.SubElement(pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar")
    bu_char.set("char", "•")
    # Set bullet size relative to text
    bu_sz = etree.SubElement(pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buSzPct")
    bu_sz.set("val", "100000")
    # Set bullet color to white
    bu_clr = etree.SubElement(pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buClr")
    srgb = etree.SubElement(bu_clr, "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
    srgb.set("val", "FFFFFF")


def add_bullet_list(slide, bullets: list[str], left: float, top: float, width: float, height: float) -> None:
    """Add a vertical bullet list with proper bullet characters."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_right = 0

    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = bullet
        style_run(run, 19)
        paragraph.level = 0
        paragraph.space_after = Pt(10)
        paragraph.line_spacing = 1.15
        _enable_bullet(paragraph)


def add_text_block(
    slide,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int,
    bold: bool = False,
    color: RGBColor = WHITE,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """Add a text block."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = box.text_frame
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    style_run(run, font_size, color=color, bold=bold)
    paragraph.alignment = align


def add_table(
    slide, columns: list[str], rows: list[list[str]], left: float, top: float, width: float, height: float
) -> None:
    """Add a styled table."""
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(columns), Inches(left), Inches(top), Inches(width), Inches(height)
    )
    table = table_shape.table
    column_width = int(Inches(width) / len(columns))
    row_height = int(Inches(height) / (len(rows) + 1))

    for index in range(len(columns)):
        table.columns[index].width = column_width
    for index in range(len(rows) + 1):
        table.rows[index].height = row_height

    for column_index, heading in enumerate(columns):
        cell = table.cell(0, column_index)
        cell.text = heading
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                style_run(run, 15, bold=True)

    for row_index, row in enumerate(rows, start=1):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT if column_index < len(columns) - 1 else PP_ALIGN.CENTER
                for run in paragraph.runs:
                    style_run(run, 14)


def build_title_slide(prs: Presentation) -> None:
    """Create slide 1."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_text_block(
        slide,
        "Agentic Azure Enterprise Landing Zone Accelerator",
        left=0.8,
        top=1.7,
        width=11.7,
        height=1.0,
        font_size=36,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text_block(
        slide,
        "AI Orchestrates · Humans Decide · Azure Executes",
        left=1.2,
        top=3.0,
        width=10.9,
        height=0.6,
        font_size=22,
        color=LIGHT_TEXT,
        align=PP_ALIGN.CENTER,
    )
    add_text_block(slide, "Microsoft Confidential", left=0.6, top=6.85, width=3.0, height=0.3, font_size=12)
    add_notes(
        slide,
        (
            "Good morning/afternoon everyone. Today I'm going to walk you through the "
            "Agentic Azure Enterprise Landing Zone Accelerator — a multi-agent AI system that "
            "transforms how organizations deploy and govern Azure Landing Zones."
        ),
    )


def build_bullet_slide(prs: Presentation, title: str, bullets: list[str], notes: str) -> None:
    """Create a standard bullet slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, title)
    add_bullet_list(slide, bullets, left=0.9, top=1.45, width=11.4, height=5.4)
    add_notes(slide, notes)


def build_architecture_slide(prs: Presentation) -> None:
    """Create slide 6 — layered architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Multi-Agent Architecture — 6 Layers")

    layers = [
        ("🖥️  Presentation", "Copilot Studio · Teams · CLI · Azure Portal"),
        ("🧠  Orchestrator", "Conductor (DAG) · Approval Gates · Session State"),
        ("🤖  Agent Layer", "14 Agents: Assessor → Scribe → Oracle → Warden → Forge → Envoy → Sentinel"),
        ("🔌  MCP Servers", "Azure Pricing (18) · Azure Platform (27) · Draw.io"),
        ("📦  IaC Modules", "Bicep (AVM) · Terraform (AVM-TF) · 8 CAF Design Areas"),
        ("☁️  Azure", "Management Groups · Subscriptions · Networking · Identity · Security"),
    ]

    y_start = 1.5
    row_height = 0.78
    for i, (label, detail) in enumerate(layers):
        y = y_start + i * row_height
        # Layer box
        box = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(3.0), Inches(0.6))
        tf = box.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = label
        style_run(run, 16, color=ACCENT_BLUE, bold=True)

        # Detail
        detail_box = slide.shapes.add_textbox(Inches(4.0), Inches(y), Inches(8.2), Inches(0.6))
        dtf = detail_box.text_frame
        drun = dtf.paragraphs[0].add_run()
        drun.text = detail
        style_run(drun, 15)

        # Arrow between layers (except last)
        if i < len(layers) - 1:
            arrow_box = slide.shapes.add_textbox(Inches(2.0), Inches(y + 0.55), Inches(0.5), Inches(0.3))
            atf = arrow_box.text_frame
            arun = atf.paragraphs[0].add_run()
            arun.text = "↓"
            style_run(arun, 12, color=LIGHT_TEXT)

    add_notes(
        slide,
        (
            "Here's the architecture in 6 layers. At the top, users interact via "
            "Copilot Studio, Teams, CLI, or Azure Portal. The Conductor orchestrates "
            "a DAG workflow with approval gates. 14 specialized agents handle specific "
            "steps — from assessment through remediation. They connect to Azure via "
            "3 MCP servers providing real-time pricing, platform operations, and "
            "diagram generation. Below that, IaC modules organized by CAF design area "
            "use Azure Verified Modules. Finally, everything deploys to Azure — "
            "management groups, subscriptions, networking, identity, and security. "
            "The Challenger agent sits adversarially across gates 1, 2, 4, and 5."
        ),
    )


def build_table_slide(
    prs: Presentation,
    title: str,
    columns: list[str],
    rows: list[list[str]],
    notes: str,
    *,
    extra_text: str | None = None,
    table_top: float = 1.45,
    table_height: float = 4.9,
) -> None:
    """Create a table-centric slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, title)
    add_table(slide, columns, rows, left=0.8, top=table_top, width=11.7, height=table_height)
    if extra_text:
        add_text_block(
            slide, extra_text, left=0.9, top=6.05, width=11.4, height=0.45, font_size=18, align=PP_ALIGN.CENTER
        )
    add_notes(slide, notes)


def build_security_slide(prs: Presentation) -> None:
    """Create slide 9."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Security Baseline — Always Enforced")
    add_bullet_list(
        slide,
        [
            "1. TLS 1.2 minimum everywhere",
            "2. HTTPS-only traffic",
            "3. No public blob access",
            "4. Managed Identity preferred (no keys/secrets)",
            "5. Azure AD-only SQL authentication",
            "6. Public network access disabled (production)",
        ],
        left=0.9,
        top=1.45,
        width=11.4,
        height=4.5,
    )
    add_text_block(
        slide,
        "Enforced at 3 layers: Code generation → Deployment preflight → Continuous monitoring",
        left=0.8,
        top=6.0,
        width=11.7,
        height=0.45,
        font_size=18,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        (
            "Security is non-negotiable in this system. We have 6 mandatory rules "
            "that are enforced at three different layers. First, at code generation "
            "— the Forge agent physically cannot produce templates that violate these "
            "rules. Second, at deployment preflight — the Challenger validates before "
            "anything hits Azure. Third, continuously — the Sentinel monitors for "
            "drift and the Mender remediates. This means even if someone manually "
            "changes a resource in the portal, the system detects and corrects it. "
            "You can't accidentally ship insecure infrastructure through this "
            "pipeline."
        ),
    )


def build_demo_slide(prs: Presentation) -> None:
    """Create slide 14."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "In Practice — From Prompt to Production")

    numbers_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.45), Inches(0.5), Inches(5.3))
    flow_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.45), Inches(10.7), Inches(5.3))
    number_frame = numbers_box.text_frame
    flow_frame = flow_box.text_frame
    flow_frame.word_wrap = True

    steps = [
        "User: 'Deploy a hub-spoke network with Azure Firewall for a SOC2-compliant workload'",
        "→ Scribe captures requirements (8 CAF areas)",
        "→ Oracle assesses WAF alignment + estimates cost",
        "→ Warden discovers existing policies, adds security baseline",
        "→ Strategist selects AVM modules, plans dependencies",
        "→ Forge generates Bicep/Terraform with full compliance",
        "→ Envoy deploys with what-if preview → User approves",
        "→ Chronicler generates TDD + runbooks",
        "→ Sentinel begins continuous monitoring",
    ]

    for index, step in enumerate(steps, start=1):
        number_paragraph = number_frame.paragraphs[0] if index == 1 else number_frame.add_paragraph()
        number_run = number_paragraph.add_run()
        number_run.text = f"{index}."
        style_run(number_run, 20, color=ACCENT_BLUE, bold=True)
        number_paragraph.space_after = Pt(12)
        number_paragraph.alignment = PP_ALIGN.RIGHT

        flow_paragraph = flow_frame.paragraphs[0] if index == 1 else flow_frame.add_paragraph()
        flow_run = flow_paragraph.add_run()
        flow_run.text = step
        style_run(flow_run, 18)
        flow_paragraph.space_after = Pt(12)
        flow_paragraph.line_spacing = 1.1

    add_notes(
        slide,
        (
            "Let me show you what this looks like in practice. A user describes their "
            "requirement — in this case, a hub-spoke network for a SOC2 workload. "
            "The pipeline automatically flows through all stages. Requirements get "
            "mapped to CAF, architecture gets WAF-assessed with real cost estimates, "
            "governance constraints are applied, code is generated using Azure "
            "Verified Modules, and deployment happens with a what-if preview. At "
            "every gate, the human approves. The whole thing generates documentation "
            "and starts continuous monitoring. What used to take weeks of architect "
            "and consultant time happens in hours with full governance throughout."
        ),
    )


def build_summary_slide(prs: Presentation) -> None:
    """Create slide 15."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Get Started")
    add_bullet_list(
        slide,
        [
            "Open source — available on GitHub",
            "Flexible — Bicep or Terraform, greenfield or brownfield",
            "Extensible — 60+ skills, custom MCP servers, modular agents",
            "Enterprise-ready — approval gates, OIDC auth, security baseline",
            "Living system — continuous monitoring + auto-remediation",
        ],
        left=0.9,
        top=1.45,
        width=11.4,
        height=3.9,
    )
    add_text_block(
        slide,
        "Try the brownfield assessment on your existing environment",
        left=1.0,
        top=5.35,
        width=11.0,
        height=0.35,
        font_size=18,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text_block(
        slide,
        "Run a pilot Landing Zone deployment",
        left=1.0,
        top=5.8,
        width=11.0,
        height=0.35,
        font_size=18,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text_block(
        slide,
        "Extend with custom skills for your organization",
        left=1.0,
        top=6.25,
        width=11.0,
        height=0.35,
        font_size=18,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        (
            "To wrap up — this accelerator is open source, flexible, and "
            "enterprise-ready. You can start with a brownfield assessment of your "
            "existing environment — no commitment, just insight. Or run a pilot "
            "deployment to see the full pipeline in action. The system is extensible "
            "— you can add custom skills, swap MCP servers, or modify agent "
            "behavior. We'd love to hear your feedback and see how you use it. "
            "Thank you."
        ),
    )


def build_deck() -> Presentation:
    """Build the presentation in memory."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_title_slide(prs)
    build_bullet_slide(
        prs,
        "About Me",
        [
            "[Your Name] — Senior Cloud Solution Architect",
            "Recently joined Microsoft",
            "Background in cloud infrastructure and enterprise architecture",
            "Passionate about AI-driven infrastructure automation",
            "Focus: Azure Landing Zones, CAF, Well-Architected Framework",
        ],
        (
            "A quick introduction — I'm [Name], a Senior Cloud Solution Architect "
            "who recently joined Microsoft. My background is in cloud "
            "infrastructure and enterprise architecture. I'm passionate about the "
            "intersection of AI and infrastructure — using intelligent automation to "
            "solve the real challenges our customers face when adopting Azure at "
            "scale. This accelerator is a project I've been building that combines "
            "that passion with Microsoft's best-practice frameworks."
        ),
    )
    build_bullet_slide(
        prs,
        "What You'll Walk Away With",
        [
            "Understand how multi-agent AI can automate Landing Zone lifecycle management",
            "See a CAF + WAF aligned approach that enforces governance from day one",
            "Learn about continuous compliance — not just deploy and forget",
            "Discover how to assess brownfield estates with AI-driven WAF analysis",
            "Get a path to accelerate your own LZ deployments from weeks to hours",
        ],
        (
            "Here's what you'll get out of the next 10 minutes. You'll understand "
            "how multi-agent AI can automate the entire Landing Zone lifecycle — "
            "from requirements through continuous monitoring. You'll see how we "
            "enforce Cloud Adoption Framework and Well-Architected Framework "
            "alignment at every stage. You'll learn about our brownfield assessment "
            "capability for existing environments. And most importantly, you'll have "
            "a clear path to accelerate your own customers' Landing Zone "
            "deployments."
        ),
    )
    build_bullet_slide(
        prs,
        "The Landing Zone Challenge",
        [
            "Landing Zone deployments take 8–16 weeks with specialized consultants",
            "Governance drift begins within days of deployment",
            "Manual IaC authoring is error-prone and inconsistent",
            "No continuous compliance monitoring post-deployment",
            "Brownfield estates are often ungoverned and undocumented",
        ],
        (
            "Let's set the context. Today, deploying a properly governed Azure "
            "Landing Zone is a significant undertaking. It requires specialized "
            "architects, weeks of planning, and expensive consultants. And here's "
            "the real problem — once it's deployed, governance drift begins almost "
            "immediately. Security baselines erode, costs creep up, and "
            "documentation goes stale. For brownfield environments — which is most "
            "enterprises — the situation is even worse. They have infrastructure "
            "that was never properly assessed against WAF or CAF. We built this "
            "accelerator to solve all of these problems with AI-driven automation "
            "that doesn't just deploy — it continuously governs."
        ),
    )
    build_bullet_slide(
        prs,
        "Why This Matters — Business Impact",
        [
            "Speed: Weeks → hours for Landing Zone deployment",
            "Consistency: Every deployment follows CAF + WAF alignment",
            "Governance: Non-negotiable security baseline enforced at every stage",
            "Cost Control: Built-in budget alerts (80/100/120%) on every deployment",
            "Continuous: Day-2 operations with auto-remediation",
            "Dual-track: Supports both Bicep and Terraform",
        ],
        (
            "The business case is compelling. We're compressing 8-16 weeks of work "
            "into hours. Every deployment is consistent because AI follows the same "
            "patterns every time — aligned to CAF design areas and WAF pillars. "
            "Security isn't optional — we enforce a 6-rule security baseline at "
            "code generation, deployment, and continuously afterward. Cost "
            "governance is built in — every deployment gets budget resources with "
            "alerts. And critically, this isn't deploy-and-forget. The system runs "
            "continuous compliance monitoring and can auto-remediate violations. "
            "Customers choose their IaC framework — Bicep or Terraform — and the "
            "system handles both."
        ),
    )
    build_architecture_slide(prs)
    build_table_slide(
        prs,
        "Specialized AI Agents",
        ["Agent", "Role", "Step"],
        SLIDE_7_ROWS,
        (
            "Each agent is purpose-built. The Scribe captures requirements mapped to "
            "all 8 CAF design areas. The Oracle runs a full WAF assessment across "
            "all 5 pillars and provides real-time cost estimates using Azure Pricing "
            "APIs. The Warden discovers existing Azure Policies and enforces our "
            "security baseline. The Forge generates IaC using Azure Verified Modules "
            "as the first choice — falling back to native resources only when no "
            "AVM exists. The Sentinel runs daily compliance scans, and the Mender "
            "auto-remediates critical violations with snapshot-based rollback. The "
            "Challenger is our adversarial reviewer — it stress-tests outputs at "
            "gates proportional to complexity."
        ),
        table_height=5.2,
    )
    build_table_slide(
        prs,
        "Full CAF Coverage",
        ["Design Area", "Coverage"],
        SLIDE_8_ROWS,
        (
            "Every module in our system maps directly to a CAF design area. This "
            "isn't accidental — it's by design. When the Scribe captures "
            "requirements, it maps them to all 8 areas. When the Oracle assesses, "
            "it evaluates against all 5 WAF pillars. When the Forge generates code, "
            "it's organized by design area. This means the output is immediately "
            "recognizable to any Azure architect — it follows the patterns they "
            "already know from the Cloud Adoption Framework."
        ),
        table_height=4.9,
    )
    build_security_slide(prs)
    build_bullet_slide(
        prs,
        "Brownfield — Assess Before You Act",
        [
            "221-check catalog synced with Azure Proactive Resiliency Library (APRL)",
            "All 5 WAF pillars: Reliability, Security, Cost, Operations, Performance",
            "Resource Graph-based discovery — no agents on VMs needed",
            "Generates current-state documentation + remediation roadmap",
            "Excel + PowerPoint reports with Microsoft branding",
        ],
        (
            "For brownfield environments, Step 0 is critical. The Assessor agent "
            "runs a comprehensive WAF assessment using a 221-check catalog aligned "
            "with the Azure Proactive Resiliency Library. It covers all 5 WAF "
            "pillars. Discovery uses Azure Resource Graph — no agents or "
            "installations required. The output includes a full current-state "
            "assessment, gap analysis, and prioritized remediation roadmap. We "
            "generate branded reports in Excel and PowerPoint format. This "
            "assessment can run standalone — you don't have to go through the full "
            "deployment pipeline. It's a great entry point for customers who want "
            "to understand their current posture before deciding on next steps."
        ),
    )
    build_table_slide(
        prs,
        "Connected to Azure — Real-Time Intelligence",
        ["MCP Server", "Capability", "Tools"],
        SLIDE_11_ROWS,
        (
            "The agents aren't working in isolation — they're connected to Azure "
            "through Model Context Protocol servers. The Azure Pricing MCP gives "
            "real-time cost estimates during architecture assessment. The Azure "
            "Platform MCP provides 27 tools including Resource Graph queries, "
            "Policy evaluation, RBAC analysis, and deployment operations. The "
            "Draw.io MCP generates architecture diagrams with proper Azure icons. "
            "This means when the Oracle says a deployment will cost a certain "
            "amount per month, that's based on actual Azure pricing, not stale "
            "estimates."
        ),
        extra_text="Agents don't guess — they query Azure in real-time",
        table_top=1.7,
        table_height=3.0,
    )
    build_bullet_slide(
        prs,
        "Beyond Deployment — Continuous Operations",
        [
            "Sentinel runs full scan daily at 06:00 UTC (compliance + drift + audit)",
            "Alert thresholds: Critical → immediate | High → 15 min | Medium → daily",
            "Mender auto-remediates critical/high with snapshot-based rollback",
            "8 built-in remediation strategies for common policy violations",
            "Escalation to human approval for medium/low severity",
        ],
        (
            "This is where we go beyond traditional accelerators. Most Landing "
            "Zone tools stop at deployment. Ours continues operating. The Sentinel "
            "agent runs a full compliance and drift scan daily. If it finds "
            "critical or high-severity violations, it immediately routes them to "
            "the Mender for auto-remediation. The Mender takes a snapshot first — "
            "so if remediation causes issues, we can roll back. For medium and low "
            "severity findings, it escalates to humans. This creates a self-healing "
            "infrastructure loop. Your Landing Zone stays compliant not just on "
            "day one, but every day after."
        ),
    )
    build_bullet_slide(
        prs,
        "Enterprise CI/CD — From Code to Cloud",
        [
            "GitHub Actions with OIDC federation — no stored credentials",
            "5-stage reusable pipeline: Resolve → Validate → Plan → Deploy → Verify",
            "Reusable workflows for consistent deployments across Landing Zones",
            "Pre-commit hooks via Lefthook: security baseline, cost governance, terraform fmt",
            "Automated assessment workflow for brownfield environments",
            "APRL sync — keeps WARA check catalog current",
        ],
        (
            "The CI/CD pipeline is fully automated with GitHub Actions using OIDC "
            "workload identity federation — no stored secrets or service principal "
            "keys. We have a 5-stage reusable pipeline that every deployment goes "
            "through: Resolve authenticates via OIDC, Validate runs security and cost "
            "checks, Plan generates what-if output, Deploy applies with timestamped "
            "names, and Verify runs compliance scans and generates the as-built TDD. "
            "Every PR goes through automated validation. The assessment workflow can "
            "run on-demand against any existing subscription. Pre-commit hooks via "
            "Lefthook catch issues before they even reach the pipeline."
        ),
    )
    build_table_slide(
        prs,
        "5-Stage Reusable Pipeline",
        ["Stage", "What Happens"],
        SLIDE_PIPELINE_ROWS,
        (
            "Every landing zone deployment — whether platform or application — goes "
            "through the same 5-stage reusable pipeline. Resolve validates OIDC auth "
            "and resolves the subscription from GitHub secrets. Validate runs security "
            "baseline, cost governance, and IaC linting — any violation fails the "
            "pipeline. Plan runs Bicep what-if or Terraform plan. Deploy applies "
            "changes with a timestamped deployment name for traceability. Verify "
            "captures a drift baseline, runs a compliance scan, and auto-generates "
            "the Technical Design Document. This DRY approach ensures consistency "
            "across all 10+ subscriptions."
        ),
        table_top=1.7,
        table_height=3.2,
        extra_text="Same pipeline for all 10 subscriptions — DRY, consistent, auditable",
    )
    build_table_slide(
        prs,
        "Landing Zone Profiles (4 Platform + 4 App)",
        ["Profile", "Purpose", "Order"],
        SLIDE_PROFILES_ROWS,
        (
            "The accelerator ships with 8 pre-built profiles — 4 platform and 4 "
            "application. Platform LZs deploy in strict order: Management first for "
            "centralized logging, then Connectivity for hub networking, Identity for "
            "AD DS and RBAC, and Security for Sentinel and Defender. Application LZs "
            "deploy in parallel after platform is ready. Online is for internet-facing "
            "apps with WAF, Corp for internal private-connectivity workloads, SAP for "
            "specialized SAP workloads with ANF and proximity groups, and Sandbox for "
            "dev/test with relaxed policies and hard budget caps. App LZ names are "
            "fully configurable via subscriptions.json — no YAML changes needed."
        ),
        table_height=4.6,
    )
    build_table_slide(
        prs,
        "Auto-Generated As-Built Documentation",
        ["TDD Section", "Contents"],
        SLIDE_TDD_ROWS,
        (
            "After every successful deployment, the pipeline automatically generates "
            "a Technical Design Document — a comprehensive as-built record of exactly "
            "what was deployed. It includes architecture diagrams with official Azure "
            "icons, a full resource inventory from Resource Graph, network topology "
            "details, security posture including all 6 baseline rules and Defender "
            "plans, live compliance scan results, cost governance configuration, and "
            "the operational model with monitoring schedules and DR plans. TDDs are "
            "uploaded as pipeline artifacts with 90-day retention. No manual "
            "documentation effort — the system documents itself."
        ),
        table_height=4.6,
        extra_text="Generated automatically in the Verify stage — zero manual effort",
    )
    build_demo_slide(prs)
    build_summary_slide(prs)
    return prs


def save_deck() -> Path:
    """Build and save the deck to disk."""
    presentation = build_deck()
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    if OUTPUT_PATH.exists():
        OUTPUT_PATH.unlink()
    presentation.save(OUTPUT_PATH)
    return OUTPUT_PATH


def main() -> None:
    """Generate the presentation and print the output path."""
    output_path = save_deck()
    print(f"Generated deck: {output_path} ({output_path.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
